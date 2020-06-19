from pdf2image import convert_from_path
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

#Cut the pdf presentation into slides
pages = convert_from_path(sys.argv[1], 500)
counter=0
paths=[]
#store slides in the directory image_temp
base = "image_temp/"+Path(sys.argv[1]).stem+"_"
Path(base).parent.mkdir(exist_ok=True)
for page in pages:
    filename=Path(base+str(counter)+".png")
    paths.append(filename)
    page.save(paths[-1].resolve(), 'PNG')
    counter+=1


#Read the empty template presentation (required for 16:9 aspect ratio), if you want 4:3 use another template file
prs = Presentation("template.pptx")
#choose a slide layout from 0-7 (Doesnt matter in this case which you take)
blank_slide_layout = prs.slide_layouts[6]   

#Place image at top left corner
left=top=Inches(0)
#height of the image
_height=Inches(7.5)
#Overwrite the first empty slide of the template presentation
pic = prs.slides[0].shapes.add_picture(str(paths[0].resolve()), left, top,height=_height)
#write all pdf slides to pptx
for path in paths[1:]:
    slide = prs.slides.add_slide(blank_slide_layout)
    pic = slide.shapes.add_picture(str(path.resolve()), left, top,height=_height)
    
#save the PPtx file
prs.save(sys.argv[2])
